#!/usr/bin/env python3
"""Deck assembly helpers for the presentation-maker skill (v2 visual layer).

The agent builds decks ONLY through these functions — never raw OOXML.
All layout constants assume a 16:9 slide (13.333 x 7.5 in).

v2 adds: a color-blocked title slide, full-bleed section/closing slides with
auto black/white text, per-slide page furniture (kicker + accent rule +
footer), rounded surface cards under stats and columns, and image slides.
Decorative shapes are named with a "deco:" prefix so validate.py can exclude
them from overlap checks.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Emu(12192000)  # 13.333 in, exact 16:9
SLIDE_H = Emu(6858000)   # 7.5 in
MARGIN = Inches(0.83)  # consistent outer margin on every slide
CONTENT_W = Inches(13.333 - 2 * 0.83)
FOOTER_ATTR = "_aionui_footer"


def _rgb(hex6):
    return RGBColor.from_string(hex6)


def _luminance(hex6):
    def channel(value):
        value /= 255.0
        return value / 12.92 if value <= 0.04045 else ((value + 0.055) / 1.055) ** 2.4

    r, g, b = (int(hex6[i:i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def _contrast(l1, l2):
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _on_color(hex6):
    """Black or white text, whichever has the higher WCAG contrast on the fill."""
    lum = _luminance(hex6)
    if _contrast(lum, _luminance("1A1A1A")) >= _contrast(lum, 1.0):
        return "1A1A1A"
    return "FFFFFF"


def _fill_slide(slide, hex6):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex6)


def _blank_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _fill_slide(slide, theme["colors"]["bg"])
    return slide


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _style_run(para, text, font_name, size, hex6, bold=False):
    para.text = text
    run = para.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(hex6)
    return run


def _set_run(para, text, theme, size, color_key, font_key="body", bold=False):
    return _style_run(para, text, theme["fonts"][font_key], size, theme["colors"][color_key], bold=bold)


def _deco_rect(slide, x, y, w, h, hex6, name, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex6)
    shape.line.fill.background()
    shape.shadow.inherit = False
    if rounded:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, TypeError):
            pass
    return shape


def _page_furniture(prs, slide, theme):
    footer = getattr(prs, FOOTER_ATTR, "") or ""
    if footer:
        _, ff = _textbox(slide, MARGIN, Inches(7.02), Inches(6.0), Inches(0.32))
        _style_run(ff.paragraphs[0], footer, theme["fonts"]["body"], 10, theme["colors"]["muted"])
    _, pf = _textbox(slide, Inches(12.1), Inches(7.02), Inches(0.9), Inches(0.32))
    pf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style_run(pf.paragraphs[0], str(len(prs.slides)), theme["fonts"]["body"], 10, theme["colors"]["muted"])


def new_deck(theme, footer=""):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    setattr(prs, FOOTER_ATTR, footer)
    return prs


def add_title_slide(prs, theme, title, subtitle="", logo_path=None):
    slide = _blank_slide(prs, theme)
    band_w = Inches(4.5)
    _deco_rect(slide, 0, 0, band_w, SLIDE_H, theme["colors"]["primary"], "deco:title-band")
    if logo_path and os.path.exists(logo_path):
        pic = slide.shapes.add_picture(logo_path, Inches(0.6), Inches(0.6), height=Inches(0.8))
        pic.name = "deco:logo"
    tx = Inches(4.5 + 0.7)
    tw = Inches(13.333 - 4.5 - 0.7 - 0.83)
    _, tf = _textbox(slide, tx, Inches(2.6), tw, Inches(1.9))
    _set_run(tf.paragraphs[0], title, theme, 44, "text", "heading", bold=True)
    if subtitle:
        _, sf = _textbox(slide, tx, Inches(4.6), tw, Inches(0.8))
        _set_run(sf.paragraphs[0], subtitle, theme, 20, "muted")


def add_section_slide(prs, theme, title, kicker=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide(slide, theme["colors"]["primary"])
    on = _on_color(theme["colors"]["primary"])
    _deco_rect(slide, MARGIN, Inches(2.3), Inches(0.9), Inches(0.06), on, "deco:rule")
    if kicker:
        _, kf = _textbox(slide, MARGIN, Inches(2.55), CONTENT_W, Inches(0.5))
        _style_run(kf.paragraphs[0], kicker.upper(), theme["fonts"]["body"], 14, on, bold=True)
    _, tf = _textbox(slide, MARGIN, Inches(3.15), CONTENT_W, Inches(1.5))
    _style_run(tf.paragraphs[0], title, theme["fonts"]["heading"], 36, on, bold=True)


def _header(prs, slide, theme, title, kicker=""):
    if kicker:
        _, kf = _textbox(slide, MARGIN, Inches(0.55), CONTENT_W, Inches(0.4))
        _set_run(kf.paragraphs[0], kicker.upper(), theme, 12, "primary", bold=True)
    _, tf = _textbox(slide, MARGIN, Inches(0.95), CONTENT_W, Inches(0.9))
    _set_run(tf.paragraphs[0], title, theme, 28, "text", "heading", bold=True)
    _deco_rect(slide, MARGIN, Inches(1.85), Inches(1.6), Inches(0.045), theme["colors"]["primary"], "deco:rule")
    _page_furniture(prs, slide, theme)


def _bullets_into(tf, theme, bullets, size=18):
    for i, item in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_run(para, str(item), theme, size, "text")
        para.space_after = Pt(10)


def add_bullets_slide(prs, theme, title, bullets, kicker=""):
    slide = _blank_slide(prs, theme)
    _header(prs, slide, theme, title, kicker)
    _, tf = _textbox(slide, MARGIN, Inches(2.1), CONTENT_W, Inches(4.6))
    _bullets_into(tf, theme, bullets)


def add_two_column_slide(prs, theme, title, left_title, left_bullets, right_title, right_bullets):
    slide = _blank_slide(prs, theme)
    _header(prs, slide, theme, title)
    col_w = Inches((13.333 - 2 * 0.83 - 0.5) / 2)
    pad = Inches(0.35)
    for idx, (col_title, col_bullets) in enumerate(((left_title, left_bullets), (right_title, right_bullets))):
        x = MARGIN + idx * (col_w + Inches(0.5))
        _deco_rect(slide, x, Inches(2.1), col_w, Inches(4.55),
                   theme["colors"]["surface"], "deco:card-%d" % idx, rounded=True)
        inner_w = col_w - pad * 2
        _, ht = _textbox(slide, x + pad, Inches(2.45), inner_w, Inches(0.5))
        _set_run(ht.paragraphs[0], col_title, theme, 18, "primary", bold=True)
        _, bt = _textbox(slide, x + pad, Inches(3.05), inner_w, Inches(3.4))
        _bullets_into(bt, theme, col_bullets, size=16)


def add_stats_slide(prs, theme, title, stats):
    if len(stats) > 4:
        raise ValueError("add_stats_slide supports at most 4 stats")
    slide = _blank_slide(prs, theme)
    _header(prs, slide, theme, title)
    n = len(stats)
    gap = Inches(0.4)
    card_w = (CONTENT_W - gap * (n - 1)) / n if n > 1 else CONTENT_W
    for i, (value, label) in enumerate(stats):
        x = MARGIN + i * (card_w + gap)
        _deco_rect(slide, x, Inches(2.5), card_w, Inches(2.6),
                   theme["colors"]["surface"], "deco:card-%d" % i, rounded=True)
        _, vf = _textbox(slide, x, Inches(3.0), card_w, Inches(1.1))
        vf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_run(vf.paragraphs[0], str(value), theme, 40, "primary", "heading", bold=True)
        _, lf = _textbox(slide, x, Inches(4.1), card_w, Inches(0.7))
        lf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_run(lf.paragraphs[0], str(label), theme, 14, "muted")


def add_quote_slide(prs, theme, quote, attribution=""):
    slide = _blank_slide(prs, theme)
    mark_box, mf = _textbox(slide, Inches(1.4), Inches(1.2), Inches(2.0), Inches(1.9))
    mark_box.name = "deco:quote-mark"
    _set_run(mf.paragraphs[0], "“", theme, 96, "primary", "heading", bold=True)
    _, tf = _textbox(slide, Inches(1.6), Inches(2.8), Inches(10.1), Inches(2.2))
    _set_run(tf.paragraphs[0], "“%s”" % quote, theme, 28, "text", "heading")
    if attribution:
        _, af = _textbox(slide, Inches(1.6), Inches(5.2), Inches(10.1), Inches(0.6))
        _set_run(af.paragraphs[0], "— %s" % attribution, theme, 16, "muted")
    _page_furniture(prs, slide, theme)


def add_image_slide(prs, theme, title, image_path, caption=""):
    slide = _blank_slide(prs, theme)
    _header(prs, slide, theme, title)
    area_x, area_y = MARGIN, Inches(2.1)
    area_w = CONTENT_W
    area_h = Inches(4.2) if caption else Inches(4.7)
    pic = slide.shapes.add_picture(image_path, area_x, area_y)
    scale = min(area_w / pic.width, area_h / pic.height)
    if scale < 1:
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
    pic.left = int(area_x + (area_w - pic.width) / 2)
    pic.top = int(area_y + (area_h - pic.height) / 2)
    if caption:
        _, cf = _textbox(slide, MARGIN, Inches(6.45), CONTENT_W, Inches(0.4))
        cf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set_run(cf.paragraphs[0], caption, theme, 14, "muted")


def add_closing_slide(prs, theme, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide(slide, theme["colors"]["primary"])
    on = _on_color(theme["colors"]["primary"])
    _, tf = _textbox(slide, MARGIN, Inches(2.7), CONTENT_W, Inches(1.2))
    _style_run(tf.paragraphs[0], title, theme["fonts"]["heading"], 40, on, bold=True)
    _, lf = _textbox(slide, MARGIN, Inches(4.1), CONTENT_W, Inches(1.5))
    for i, item in enumerate(lines):
        para = lf.paragraphs[0] if i == 0 else lf.add_paragraph()
        _style_run(para, str(item), theme["fonts"]["body"], 16, on)
        para.space_after = Pt(10)


def save_deck(prs, path):
    path = os.path.abspath(path)
    prs.save(path)
    return path

