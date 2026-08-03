#!/usr/bin/env python3
"""Geometric QA for generated .pptx decks.

Heuristics (no rendering deps):
- overflow: estimated wrapped text height exceeds its frame height
- overlap: two visible shapes intersect by >15% of the smaller one
- off_slide: shape extends beyond slide bounds

Usage: validate.py <deck.pptx> | validate.py --self-test
"""
import json
import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

SLIDE_W = 12192000  # EMU, 13.333 in
SLIDE_H = 6858000   # EMU, 7.5 in
EMU_PER_PT = 12700
CHAR_WIDTH_FACTOR = 0.55   # avg glyph width as fraction of font size
LINE_HEIGHT_FACTOR = 1.35  # line height as fraction of font size
OVERLAP_TOLERANCE = 0.15


def _font_size_pt(para, default=18.0):
    for run in para.runs:
        if run.font.size is not None:
            return run.font.size.pt
    return default


def _estimate_text_height_emu(shape):
    tf = shape.text_frame
    total_pt = 0.0
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs)
        size = _font_size_pt(para)
        if not text.strip():
            total_pt += size * LINE_HEIGHT_FACTOR
            continue
        width_pt = shape.width / EMU_PER_PT
        chars_per_line = max(1, int(width_pt / (size * CHAR_WIDTH_FACTOR)))
        lines = max(1, -(-len(text) // chars_per_line))  # ceil div
        total_pt += lines * size * LINE_HEIGHT_FACTOR
        if para.space_after is not None:
            total_pt += para.space_after.pt
    return int(total_pt * EMU_PER_PT)


def _rect(shape):
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def _is_background_sized(shape):
    return (shape.width or 0) * (shape.height or 0) >= 0.9 * SLIDE_W * SLIDE_H


def _is_decorative(shape):
    """Helpers name intentional background shapes (cards, bands, rules)
    with a 'deco:' prefix — they may sit under content by design."""
    return (shape.name or "").startswith("deco:")


def _shape_name(shape):
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text.strip().splitlines()[0][:30] if shape.text_frame.text.strip() else ""
    return "%s(%r)" % (shape.shape_type, text) if text else str(shape.shape_type)


def validate(path):
    prs = Presentation(path)
    issues = []
    for i, slide in enumerate(prs.slides, start=1):
        shapes = [s for s in slide.shapes if s.width and s.height]
        for shape in shapes:
            l, t, r, b = _rect(shape)
            if l < 0 or t < 0 or r > SLIDE_W or b > SLIDE_H:
                issues.append({"slide": i, "shape": _shape_name(shape), "type": "off_slide",
                               "detail": "bounds (%d,%d,%d,%d) exceed slide" % (l, t, r, b)})
            if shape.has_text_frame and shape.text_frame.text.strip():
                est = _estimate_text_height_emu(shape)
                if est > shape.height:
                    issues.append({"slide": i, "shape": _shape_name(shape), "type": "overflow",
                                   "detail": "text needs ~%.2fin, frame is %.2fin"
                                             % (est / 914400.0, shape.height / 914400.0)})
        visible = [s for s in shapes if not _is_background_sized(s) and not _is_decorative(s)]
        for a_idx in range(len(visible)):
            for b_idx in range(a_idx + 1, len(visible)):
                a, b = visible[a_idx], visible[b_idx]
                al, at, ar, ab = _rect(a)
                bl, bt, br, bb = _rect(b)
                iw = min(ar, br) - max(al, bl)
                ih = min(ab, bb) - max(at, bt)
                if iw > 0 and ih > 0:
                    smaller = min((ar - al) * (ab - at), (br - bl) * (bb - bt))
                    if smaller and (iw * ih) / smaller > OVERLAP_TOLERANCE:
                        issues.append({"slide": i, "shape": _shape_name(a), "type": "overlap",
                                       "detail": "overlaps %s by %d%%"
                                                 % (_shape_name(b), 100 * iw * ih // smaller)})
    return {"slides": len(prs.slides), "ok": not issues, "issues": issues}


def _build_fixture(kind):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def box(x, y, w, h, text, size):
        b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        b.text_frame.word_wrap = True
        b.text_frame.paragraphs[0].text = text
        b.text_frame.paragraphs[0].runs[0].font.size = Pt(size)
        return b

    if kind == "broken":
        box(0.8, 0.8, 3.0, 0.4, "word " * 60, 24)          # overflow
        box(1.0, 3.0, 4.0, 1.0, "left card", 18)           # overlap pair
        box(2.0, 3.2, 4.0, 1.0, "right card", 18)
        box(11.0, 6.8, 4.0, 1.5, "runs off the slide", 18)  # off_slide
    elif kind == "cards":
        # A decorative surface card with content text on top of it must NOT
        # be flagged as overlap — this is the v2 card pattern.
        from pptx.enum.shapes import MSO_SHAPE
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(4.0), Inches(2.5))
        card.name = "deco:card-0"
        box(1.2, 2.4, 3.6, 0.8, "42%", 40)
        box(1.2, 3.4, 3.6, 0.6, "retention", 14)
    elif kind == "table":
        # An in-bounds native table under a title must pass untouched.
        frame = slide.shapes.add_table(3, 3, Inches(0.83), Inches(2.05),
                                       Inches(11.67), Inches(2.0))
        frame.table.cell(0, 0).text = "Benefit"
        box(0.8, 0.95, 11.7, 0.9, "Table slide title", 28)
    elif kind == "table-offslide":
        slide.shapes.add_table(2, 2, Inches(11.0), Inches(6.5),
                               Inches(4.0), Inches(2.0))
    else:  # clean
        box(0.8, 0.8, 11.7, 1.0, "Clean title", 28)
        box(0.8, 2.2, 11.7, 3.0, "One short line of body text.", 18)

    f = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(f.name)
    return f.name


def self_test():
    paths = []
    try:
        broken_path = _build_fixture("broken")
        paths.append(broken_path)
        broken = validate(broken_path)
        kinds = {issue["type"] for issue in broken["issues"]}
        assert not broken["ok"], "broken fixture passed"
        for expected in ("overflow", "overlap", "off_slide"):
            assert expected in kinds, "missing %s in %s" % (expected, sorted(kinds))
        clean_path = _build_fixture("clean")
        paths.append(clean_path)
        clean = validate(clean_path)
        assert clean["ok"], "clean fixture flagged: %s" % clean["issues"]
        cards_path = _build_fixture("cards")
        paths.append(cards_path)
        cards = validate(cards_path)
        assert cards["ok"], "card fixture flagged: %s" % cards["issues"]
        table_path = _build_fixture("table")
        paths.append(table_path)
        table_ok = validate(table_path)
        assert table_ok["ok"], "table fixture flagged: %s" % table_ok["issues"]
        table_off_path = _build_fixture("table-offslide")
        paths.append(table_off_path)
        table_off = validate(table_off_path)
        off_kinds = {issue["type"] for issue in table_off["issues"]}
        assert not table_off["ok"] and "off_slide" in off_kinds, \
            "off-slide table not flagged: %s" % table_off["issues"]
    finally:
        for path in paths:
            try:
                os.unlink(path)
            except OSError:
                pass
    print("SELF-TEST OK")
    return 0


def main():
    if len(sys.argv) != 2:
        print(__doc__)
        return 2
    if sys.argv[1] == "--self-test":
        return self_test()
    report = validate(sys.argv[1])
    print(json.dumps(report, indent=2))
    return 0 if report["ok"] else 1


if __name__ == "__main__":
    sys.exit(main())
